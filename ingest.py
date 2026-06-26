"""
ingest.py – Document parsing, chunking, and LanceDB ingestion logic.

Dual-routing strategy
─────────────────────
• Standard documents  (PDF, TXT, PPTX) → Recursive Character Text Splitter
• Tabular  documents  (CSV, XLSX)       → Entity-Anchored Serialisation via pandas

Both routes produce a list of plain strings ("chunks") plus metadata dicts.
The caller (main.py) generates embeddings and writes everything to LanceDB.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any, Dict, Generator, List, Tuple

import pandas as pd

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Type aliases
# ---------------------------------------------------------------------------

Chunk = str
Metadata = Dict[str, Any]
ChunkWithMeta = Tuple[Chunk, Metadata]


# ===========================================================================
# ROUTE A – Standard text splitter (PDF / TXT / PPTX)
# ===========================================================================

def _recursive_split(text: str, chunk_size: int, overlap: int) -> List[str]:
    """
    A lightweight recursive character text splitter that mirrors the behaviour
    of LangChain's RecursiveCharacterTextSplitter without the dependency.

    Priority of split separators:
      1. Paragraph boundaries  (\n\n)
      2. Single newlines        (\n)
      3. Sentences              (. / ! / ?)
      4. Words                  (space)
      5. Characters             (hard cut)
    """
    separators = ["\n\n", "\n", ". ", "! ", "? ", " ", ""]
    return _split_with_separators(text, separators, chunk_size, overlap)


def _split_with_separators(
    text: str,
    separators: List[str],
    chunk_size: int,
    overlap: int,
) -> List[str]:
    """Recursively split `text` using the first separator that produces
    sub-strings short enough to fit within `chunk_size`."""

    if len(text) <= chunk_size:
        return [text.strip()] if text.strip() else []

    # Try each separator in priority order
    for sep in separators:
        if sep == "":
            # Hard character cut – last resort
            return _hard_cut(text, chunk_size, overlap)

        parts = text.split(sep)
        if len(parts) > 1:
            chunks: List[str] = []
            current = ""
            for part in parts:
                candidate = (current + sep + part).strip() if current else part.strip()
                if len(candidate) <= chunk_size:
                    current = candidate
                else:
                    if current:
                        chunks.append(current)
                    # If a single part is itself larger than chunk_size, recurse
                    if len(part) > chunk_size:
                        sub_chunks = _split_with_separators(
                            part, separators[separators.index(sep) + 1 :], chunk_size, overlap
                        )
                        chunks.extend(sub_chunks)
                        current = ""
                    else:
                        current = part.strip()
            if current:
                chunks.append(current)

            # Apply overlap: prepend tail of previous chunk to next chunk
            return _apply_overlap(chunks, overlap)

    return [text.strip()]


def _hard_cut(text: str, chunk_size: int, overlap: int) -> List[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return chunks


def _apply_overlap(chunks: List[str], overlap: int) -> List[str]:
    """Prepend the last `overlap` characters of chunk[i-1] to chunk[i]."""
    if overlap <= 0 or len(chunks) <= 1:
        return chunks
    result = [chunks[0]]
    for i in range(1, len(chunks)):
        tail = result[-1][-overlap:] if len(result[-1]) >= overlap else result[-1]
        result.append((tail + " " + chunks[i]).strip())
    return result


# ---------------------------------------------------------------------------
# PDF parser
# ---------------------------------------------------------------------------

def _parse_pdf(path: Path) -> str:
    """Extract all text from a PDF file using pypdf."""
    try:
        from pypdf import PdfReader  # lazy import – keeps startup fast
    except ImportError as exc:
        raise RuntimeError("pypdf is not installed. Run: pip install pypdf") from exc

    reader = PdfReader(str(path))
    pages: List[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text)
    return "\n\n".join(pages)


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------

def _parse_pptx(path: Path) -> str:
    """Extract text from all slides of a PowerPoint file."""
    try:
        from pptx import Presentation  # lazy import
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is not installed. Run: pip install python-pptx"
        ) from exc

    prs = Presentation(str(path))
    slide_texts: List[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slide_texts.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
    return "\n\n".join(slide_texts)


# ---------------------------------------------------------------------------
# Standard-document route entry point
# ---------------------------------------------------------------------------

def parse_standard_document(
    path: Path,
    chunk_size: int,
    chunk_overlap: int,
) -> List[ChunkWithMeta]:
    """
    Parse a PDF, TXT, or PPTX file, split into overlapping chunks, and
    attach metadata to each chunk.

    Returns:
        List of (chunk_text, metadata_dict) tuples.
    """
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        raw_text = _parse_pdf(path)
    elif suffix == ".txt":
        raw_text = path.read_text(encoding="utf-8", errors="replace")
    elif suffix in {".ppt", ".pptx"}:
        raw_text = _parse_pptx(path)
    else:
        raise ValueError(f"Unsupported standard document type: {suffix}")

    raw_text = raw_text.strip()
    if not raw_text:
        logger.warning("No text extracted from %s", path.name)
        return []

    chunks = _recursive_split(raw_text, chunk_size, chunk_overlap)

    results: List[ChunkWithMeta] = []
    for idx, chunk in enumerate(chunks):
        if not chunk.strip():
            continue
        meta: Metadata = {
            "source_file": path.name,
            "source_type": "standard",
            "chunk_index": idx,
        }
        results.append((chunk, meta))

    logger.info("Standard parse: %d chunks from '%s'", len(results), path.name)
    return results


# ===========================================================================
# ROUTE B – Entity-Anchored Serialisation (CSV / XLSX)
# ===========================================================================

def _serialise_row(entity: str, row: pd.Series, entity_col: str) -> str:
    """
    Build a semantically dense string for one dataframe row.

    The primary entity name (first column value) is repeated next to every
    fact so that BM25 keyword matching can locate it even when the query
    references the entity alongside a specific attribute.

    Example output (entity = "Alice Johnson"):
        "Regarding Alice Johnson: The Age for Alice Johnson is 34.
         The Department for Alice Johnson is Engineering.
         The Salary for Alice Johnson is 95000."
    """
    facts: List[str] = []
    for col, val in row.items():
        if col == entity_col:
            continue  # entity column itself is already in the header
        val_str = str(val).strip()
        if val_str and val_str.lower() not in {"nan", "none", ""}:
            facts.append(f"The {col} for {entity} is {val_str}.")

    if not facts:
        return f"Regarding {entity}: (no additional data)"

    facts_text = " ".join(facts)
    return f"Regarding {entity}: {facts_text}"


def parse_tabular_document(path: Path) -> List[ChunkWithMeta]:
    """
    Read a CSV or XLSX file with pandas and apply Entity-Anchored Serialisation.

    Each row becomes one "chunk" – a self-contained semantic string that
    repeats the primary entity name (first column) next to every fact.
    This maximises BM25 recall for queries like
    "What is Alice Johnson's salary?" even without semantic similarity.

    Returns:
        List of (serialised_row_string, metadata_dict) tuples.
    """
    suffix = path.suffix.lower()

    if suffix == ".csv":
        df = pd.read_csv(path, dtype=str)
    elif suffix in {".xls", ".xlsx"}:
        df = pd.read_excel(path, dtype=str, engine="openpyxl")
    else:
        raise ValueError(f"Unsupported tabular document type: {suffix}")

    # Drop completely empty rows
    df.dropna(how="all", inplace=True)
    df.reset_index(drop=True, inplace=True)

    if df.empty:
        logger.warning("No data rows found in '%s'", path.name)
        return []

    # The first column is treated as the primary entity identifier
    entity_col: str = df.columns[0]
    logger.info(
        "Tabular parse: using column '%s' as primary entity for '%s'",
        entity_col,
        path.name,
    )

    results: List[ChunkWithMeta] = []
    for row_idx, row in df.iterrows():
        entity = str(row[entity_col]).strip()
        if not entity or entity.lower() in {"nan", "none"}:
            entity = f"Row {row_idx}"

        chunk = _serialise_row(entity, row, entity_col)

        meta: Metadata = {
            "source_file": path.name,
            "source_type": "tabular",
            "chunk_index": int(str(row_idx)),
            "entity": entity,
            "entity_column": entity_col,
        }
        results.append((chunk, meta))

    logger.info("Tabular parse: %d rows serialised from '%s'", len(results), path.name)
    return results


# ===========================================================================
# Public router function
# ===========================================================================

def route_and_parse(
    path: Path,
    chunk_size: int,
    chunk_overlap: int,
) -> List[ChunkWithMeta]:
    """
    Inspect the file extension and dispatch to the correct parsing route.

    Args:
        path:          Absolute path to the uploaded file.
        chunk_size:    Target chunk size in characters (standard docs only).
        chunk_overlap: Overlap in characters between adjacent chunks (std only).

    Returns:
        List of (chunk_text, metadata_dict) ready for embedding and ingestion.

    Raises:
        ValueError: If the file extension is not supported.
    """
    suffix = path.suffix.lower()

    STANDARD_TYPES = {".pdf", ".txt", ".ppt", ".pptx"}
    TABULAR_TYPES = {".csv", ".xls", ".xlsx"}

    if suffix in STANDARD_TYPES:
        return parse_standard_document(path, chunk_size, chunk_overlap)
    elif suffix in TABULAR_TYPES:
        return parse_tabular_document(path)
    else:
        supported = sorted(STANDARD_TYPES | TABULAR_TYPES)
        raise ValueError(
            f"Unsupported file type '{suffix}'. Supported: {supported}"
        )
